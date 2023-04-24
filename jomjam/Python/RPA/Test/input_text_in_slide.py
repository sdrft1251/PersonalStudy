import win32com.client
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def text_on_shape(shape, input_text, font_name = "@@", font_size=10, bold=False):
    text_frame = shape.TextFrame
    text_frame.DeleteText

    text_frame.TextRange.Text = input_text
    text_frame.TextRange.Font.Name = font_name
    text_frame.TextRange.Font.Size = font_size
    text_frame.TextRange.Font.Bold = bold



Application = win32com.client.Dispatch("PowerPoint.Application")
Application.Visible = True

file_path = "@@"

ppt_file = Application.Presentations.Open(file_path, ReadOnly= False)

input_info_data={
    1 : ["@@","@@"],
    2 : ["@@","@@"],
    3 : ["@@","@@"],
    4 : ["@@","@@"],
    5 : ["@@","@@"],
    6 : ["X","O"]
    }
for slide in ppt_file.Slides:
    for shape in slide.Shapes:
        print(shape.name)
        print(type(shape))
        if "Table_0" in shape.name:
            for ir in range(2,4):
                for ic in range(1,7):
                    shape.Table.Cell(ir,ic).shape.TextFrame.TextRange.Text = input_info_data[ic][ir-2]
                    shape.Table.Cell(ir,ic).shape.TextFrame.TextRange.Font.Name = "@@"
                    shape.Table.Cell(ir,ic).shape.TextFrame.TextRange.Font.Size = 20
                    shape.Table.Cell(ir,ic).shape.TextFrame.TextRange.Font.Bold = True

        elif shape.name == "Title_0":
            shape.TextFrame.TextRange.Font.Name = "@@"
            shape.TextFrame.TextRange.Font.Size = 30
            shape.TextFrame.TextRange.Font.Bold = True
            shape.TextFrame.TextRange.Text = "새로운 제목 기입"
        else:
            shape.TextFrame.TextRange.Text = "띄어쓰기 및 단락 구분 가능?\n줄바꿈\n\t탭시도\n 띄기시도"
            shape.TextFrame.TextRange.Font.Name = "@@"
            shape.TextFrame.TextRange.Font.Size = 25

