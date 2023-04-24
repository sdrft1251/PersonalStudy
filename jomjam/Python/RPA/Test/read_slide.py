import win32com.client
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def text_on_shape(shape, input_text, font_size=10, bold=False):
    text_frame = shape.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]

    p.alighnment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = input.text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.name = None


Application = win32com.client.Dispatch("PowerPoint.Application")
Application.Visible = True

file_path = "C:\\Users\\7039966\\Desktop\\test_file.pptx"

ppt_file = Application.Presentations.Open(file_path, ReadOnly= False)

for slide in ppt_file.Slides:
    for shape in slide.Shapes:
        print(shape.name)
        print(type(shape))
        if "Table" in shape.name:
            for ic in range(shape.Table.Columns.Count):
                for ir in range(shape.Table.Rows.Count):
                    print(shape.Table.Cell(ir+1,ic+1).shape.TextFrame.TextRange.Text)
            #print(shape.Table.Cell(2,2).shape.TextFrame.TextRange.Text)
            #for cell_shape in shape.Table.Cell:
            #    print(cell_shape.TextFrame.TextRange.Text)
        elif shape.name == "edit_0":
            shape.TextFrame.TextRange.Text = "Edit Complete! \n Enter!!"
        else:

            for lin in shape.TextFrame2.TextRange.Lines:
                print(lin.Text)
            #tmp_i = 1
            #while(1):
            #    text_data = shape.TextFrame.TextRange.Lines(tmp_i,tmp_i).Text
            #    if len(text_data) == 0:
            #        break
            #    print(text_data)
            #    tmp_i += 1

